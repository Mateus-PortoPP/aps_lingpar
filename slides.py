from pptx import Presentation

prs = Presentation()

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "FisioLang"
slide.placeholders[1].text = "Linguagem de Programação para Fisioterapeutas"

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Motivação"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Por que criar a FisioLang?"
for item in [
    "Auxiliar fisioterapeutas no planejamento programático de sessões",
    "Registrar dados clínicos de forma consistente",
    "Integrar dados reais de sensores e datasets"
]:
    p = tf.add_paragraph()
    p.text = f"- {item}"
    p.level = 1

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Características Principais"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Funcionalidades da FisioLang"
for item in [
    "Variáveis por atribuição implícita",
    "Condicionais: se / caso contrario",
    "Loops: treino e praticar_ate",
    "I/O e registros: resultado(), pausa(), registrar()"
]:
    p = tf.add_paragraph()
    p.text = f"- {item}"
    p.level = 1

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Curiosidades"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Diferenciais da Linguagem"
for item in [
    "Temática focada em fisioterapia",
    "Funções integradas: melhora(), ler_sensor(), calcular_consultas()",
    "Importação de CSV via importar_dados()"
]:
    p = tf.add_paragraph()
    p.text = f"- {item}"
    p.level = 1

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Exemplo de Programa"
code = (
    "sessao {\n"
    "  paciente { peso=75; altura=170; lesao=\"joelho\"; sessoes=5; }\n"
    "  importar_dados(\"physical_therapy_exercises.csv\");\n"
    "  tempo = 30;\n"
    "  resultado(tempo);\n"
    "  se melhorou(ler_sensor(\"amplitude\") > 80) {\n"
    "    tempo = tempo - 5;\n"
    "  } caso contrario {\n"
    "    tempo = tempo + 10;\n"
    "  }\n"
    "  treino(tempo > 0) { resultado(tempo); tempo = tempo - 1; }\n"
    "  praticar_ate(melhora() >= 80) { registrar(ler_sensor(\"forca\"), tempo); }\n"
    "  consultas = calcular_consultas(peso, altura, sessoes);\n"
    "  resultado(consultas);\n"
    "}\n"
    "fim_sessao"
)
tf = slide.shapes.placeholders[1].text_frame
tf.text = code

prs.save('FisioLang_Presentation.pptx')
print("Apresentação salva em FisioLang_Presentation.pptx")
